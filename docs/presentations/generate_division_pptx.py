#!/usr/bin/env python3
"""
Division-Level Access PowerPoint Generator
Creates a professional PowerPoint presentation for the Division-Level Access feature.

Requirements:
    pip install python-pptx

Usage:
    python generate_division_pptx.py

Output:
    Division-Level-Access-Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors (using RGBColor instead of RGBColor)
BLUE_PRIMARY = RGBColor(30, 64, 175)      # #1E40AF
BLUE_SECONDARY = RGBColor(59, 130, 246)   # #3B82F6
GREEN_SUCCESS = RGBColor(34, 197, 94)     # #22C55E
GRAY_TEXT = RGBColor(71, 85, 105)         # #475569
GRAY_LIGHT = RGBColor(241, 245, 249)      # #F1F5F9
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE_PRIMARY
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(191, 219, 254)  # Light blue
    p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "EduSync School Information System | December 2025"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(147, 197, 253)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullet_points):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE_PRIMARY
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content area
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if point.startswith('##'):
            # Subheading
            p.text = point.replace('##', '').strip()
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = BLUE_PRIMARY
            p.space_before = Pt(20)
        elif point.startswith('-'):
            # Bullet point
            p.text = "• " + point.replace('-', '').strip()
            p.font.size = Pt(18)
            p.font.color.rgb = GRAY_TEXT
            p.space_before = Pt(8)
            p.level = 1
        else:
            # Normal text
            p.text = point
            p.font.size = Pt(20)
            p.font.color.rgb = GRAY_TEXT
            p.space_before = Pt(12)
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE_PRIMARY
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Table
    cols = len(headers)
    table_rows = len(rows) + 1  # +1 for header
    
    table = slide.shapes.add_table(
        table_rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5 * table_rows)
    ).table
    
    # Header row
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_SECONDARY
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = GRAY_TEXT
            p.alignment = PP_ALIGN.CENTER
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GRAY_LIGHT
    
    return slide

def add_dashboard_slide(prs, title, subtitle):
    """Add a professional dashboard mockup slide with visual cards"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE_PRIMARY
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Subtitle/Description
    if subtitle:
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(9), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY_TEXT
    
    # Dashboard container background
    container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.85), Inches(9.4), Inches(5.3)
    )
    container.fill.solid()
    container.fill.fore_color.rgb = RGBColor(248, 250, 252)
    container.line.color.rgb = RGBColor(226, 232, 240)
    
    # Division header bar inside container
    div_header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.05), Inches(9), Inches(0.6)
    )
    div_header.fill.solid()
    div_header.fill.fore_color.rgb = WHITE
    div_header.line.color.rgb = RGBColor(226, 232, 240)
    
    # Division name
    div_text = slide.shapes.add_textbox(Inches(0.7), Inches(2.15), Inches(5), Inches(0.4))
    tf = div_text.text_frame
    p = tf.paragraphs[0]
    p.text = "📍 Division of City Schools - Manila"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_PRIMARY
    
    # User role
    user_text = slide.shapes.add_textbox(Inches(7), Inches(2.15), Inches(2.3), Inches(0.4))
    tf = user_text.text_frame
    p = tf.paragraphs[0]
    p.text = "👤 Superintendent"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY_TEXT
    p.alignment = PP_ALIGN.RIGHT
    
    # Statistics cards
    stats = [
        ("🏫", "42", "Schools", BLUE_PRIMARY),
        ("👨‍🎓", "25,432", "Students", RGBColor(16, 185, 129)),
        ("👨‍🏫", "1,245", "Teachers", RGBColor(139, 92, 246)),
        ("📈", "+3.2%", "Growth", RGBColor(245, 158, 11)),
    ]
    
    card_width = Inches(2.05)
    card_height = Inches(1.2)
    start_x = Inches(0.5)
    gap = Inches(0.2)
    card_y = Inches(2.85)
    
    for i, (icon, value, label, color) in enumerate(stats):
        x = start_x + (card_width + gap) * i
        
        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, card_y, card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(226, 232, 240)
        
        # Icon and value
        val_box = slide.shapes.add_textbox(x, card_y + Inches(0.15), card_width, Inches(0.6))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{icon} {value}"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        lbl_box = slide.shapes.add_textbox(x, card_y + Inches(0.7), card_width, Inches(0.35))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY_TEXT
        p.alignment = PP_ALIGN.CENTER
    
    # Schools Overview Section
    section_y = Inches(4.25)
    
    # Section header
    section_header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), section_y, Inches(9), Inches(0.5)
    )
    section_header.fill.solid()
    section_header.fill.fore_color.rgb = WHITE
    section_header.line.color.rgb = RGBColor(226, 232, 240)
    
    section_title = slide.shapes.add_textbox(Inches(0.7), section_y + Inches(0.1), Inches(4), Inches(0.35))
    tf = section_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🏫 Schools Overview"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_PRIMARY
    
    view_all = slide.shapes.add_textbox(Inches(7.5), section_y + Inches(0.1), Inches(1.8), Inches(0.35))
    tf = view_all.text_frame
    p = tf.paragraphs[0]
    p.text = "View All →"
    p.font.size = Pt(12)
    p.font.color.rgb = BLUE_SECONDARY
    p.alignment = PP_ALIGN.RIGHT
    
    # School cards
    schools = [
        ("Rizal Elem.", "523"),
        ("Bonifacio ES", "412"),
        ("Mabini Central", "687"),
        ("Aguinaldo ES", "345"),
    ]
    
    school_card_width = Inches(2.05)
    school_card_height = Inches(0.9)
    school_y = section_y + Inches(0.65)
    
    for i, (school_name, students) in enumerate(schools):
        x = start_x + (school_card_width + gap) * i
        
        # School card
        school_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, school_y, school_card_width, school_card_height
        )
        school_card.fill.solid()
        school_card.fill.fore_color.rgb = WHITE
        school_card.line.color.rgb = RGBColor(226, 232, 240)
        
        # School name
        name_box = slide.shapes.add_textbox(x, school_y + Inches(0.1), school_card_width, Inches(0.35))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"🏫 {school_name}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = GRAY_TEXT
        p.alignment = PP_ALIGN.CENTER
        
        # Student count
        count_box = slide.shapes.add_textbox(x, school_y + Inches(0.45), school_card_width, Inches(0.35))
        tf = count_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"👨‍🎓 {students} students"
        p.font.size = Pt(11)
        p.font.color.rgb = BLUE_SECONDARY
        p.alignment = PP_ALIGN.CENTER
    
    # More indicator
    more_box = slide.shapes.add_textbox(Inches(8.7), school_y + Inches(0.3), Inches(0.6), Inches(0.35))
    tf = more_box.text_frame
    p = tf.paragraphs[0]
    p.text = "..."
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GRAY_TEXT
    
    # Footer note
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(9), Inches(0.3))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Real-time dashboard with server-side aggregation for optimal performance"
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = RGBColor(148, 163, 184)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_stats_slide(prs, title, stats):
    """Add a slide with statistics cards"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE_PRIMARY
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Stats cards (4 across)
    card_width = Inches(2.1)
    card_height = Inches(1.5)
    start_x = Inches(0.5)
    gap = Inches(0.2)
    
    for i, (label, value, color) in enumerate(stats[:4]):
        x = start_x + (card_width + gap) * i
        
        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(226, 232, 240)
        
        # Value
        val_box = slide.shapes.add_textbox(x, Inches(1.6), card_width, Inches(0.8))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(value)
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color if color else BLUE_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        lbl_box = slide.shapes.add_textbox(x, Inches(2.4), card_width, Inches(0.4))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY_TEXT
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(
        prs,
        "🏢 Division-Level Access",
        "A New Era of Consolidated School Management"
    )
    
    # Slide 2: The Challenge
    add_content_slide(prs, "📋 The Challenge", [
        "## Current Pain Points for Division Offices:",
        "- Multiple logins required - logging into each school's system individually",
        "- Manual aggregation - compiling reports from 20-50+ schools by hand",
        "- Delayed reporting - weeks to consolidate SF5/SF6/SF7 reports",
        "- No real-time data - cannot see current enrollment or attendance",
        "- Inconsistent formats - each school submits reports differently",
        "",
        "Division personnel spend 40+ hours per month just collecting and consolidating school data."
    ])
    
    # Slide 3: The Solution
    add_content_slide(prs, "✨ The Solution: One Login, All Schools", [
        "## What is Division-Level Access?",
        "A centralized portal that gives Division Office personnel the ability to:",
        "- View all schools in their jurisdiction from one dashboard",
        "- Access real-time data across all schools instantly",
        "- Generate consolidated reports (SF5, SF6, SF7) in seconds",
        "- Compare school performance side-by-side",
        "- Filter by district or school for focused analysis",
        "- Export data in PDF, CSV, or Excel formats",
        "",
        "\"One Division, One View, Complete Oversight\""
    ])
    
    # Slide 4: User Roles
    add_table_slide(prs, "👥 User Roles & Permissions", 
        ["Role", "Description", "Access Level"],
        [
            ["🎓 Superintendent", "Division head", "Full access"],
            ["👔 Supervisor", "Division supervisor", "View all, limited admin"],
            ["📊 Analyst", "Data specialist", "Reports only"],
            ["👁️ Viewer", "Basic access", "Dashboard only"],
            ["📍 PSDS", "District supervisor", "Assigned district(s)"],
        ]
    )
    
    # Slide 5: Dashboard (with visual mockup)
    add_dashboard_slide(prs, "📊 Division Dashboard", 
        "Your Command Center for All Schools"
    )
    
    # Slide 6: Cascading Filters
    add_content_slide(prs, "🔍 Cascading Filters", [
        "## Smart Filtering: Division → District → School",
        "",
        "- Select a District → School dropdown filters to that district only",
        "- Select a School → All data filters to that specific school",
        "- Clear Filters → Returns to division-wide view",
        "",
        "## Benefits:",
        "- 🎯 Focus Analysis - Drill down to specific areas",
        "- ⚡ Fast Navigation - No need to switch pages",
        "- 💾 Persistent Selection - Filters are remembered across sessions"
    ])
    
    # Slide 7: DepEd Reports
    add_table_slide(prs, "📋 DepEd Reports Integration",
        ["Form", "Name", "Description", "Status"],
        [
            ["SF5", "Report on Promotion", "End-of-year statistics", "✅ Ready"],
            ["SF6", "Summarized Promotion", "Division-wide summary", "✅ Ready"],
            ["SF7", "Personnel Assignment", "All personnel list", "✅ Ready"],
        ]
    )
    
    # Slide 8: Time Savings
    add_table_slide(prs, "⏱️ Time Savings",
        ["Task", "Before", "After", "Savings"],
        [
            ["Compile SF5 (42 schools)", "3 days", "5 minutes", "99%"],
            ["Generate SF6 summary", "1 day", "2 minutes", "99%"],
            ["Collect SF7 personnel", "2 weeks", "10 minutes", "99%"],
            ["Dashboard data update", "Weekly", "Real-time", "100%"],
        ]
    )
    
    # Slide 9: Performance
    add_table_slide(prs, "⚡ Performance Optimizations",
        ["Feature", "Before", "After", "Improvement"],
        [
            ["Dashboard Load", "3.5s", "<1s", "71% faster"],
            ["Schools Grid (50 schools)", "5s", "<1s", "80% faster"],
            ["SF5 Report Generation", "8s", "<2s", "75% faster"],
            ["Personnel Summary", "4s", "<0.5s", "87% faster"],
        ]
    )
    
    # Slide 10: Security
    add_content_slide(prs, "🔒 Security & Audit", [
        "## Enterprise-Grade Security:",
        "- 🔐 Role-Based Access - Users see only what they're authorized to",
        "- 🔑 Secure Authentication - Firebase Authentication with 2FA support",
        "- 📝 Audit Logging - Every action is tracked and recorded",
        "- 🔒 Data Encryption - All data encrypted at rest and in transit",
        "- ⏰ Session Management - Auto-logout after inactivity",
        "",
        "## Compliance:",
        "- RA 10173 (Data Privacy Act) compliant",
        "- DepEd data handling guidelines",
        "- Government security standards"
    ])
    
    # Slide 11: Benefits
    add_content_slide(prs, "🎯 Benefits Summary", [
        "## For Division Personnel:",
        "- ⏰ Time Savings - 40+ hours/month saved on data collection",
        "- 📊 Real-Time Data - Always up-to-date information",
        "- 🎯 Better Decisions - Data-driven insights across schools",
        "- 📋 Compliance - Faster DepEd report submission",
        "",
        "## For Schools:",
        "- No manual submission - data automatically available",
        "- Less follow-ups from division office",
        "- Reduced phone inquiries for reports"
    ])
    
    # Slide 12: Thank You
    add_title_slide(
        prs,
        "🙏 Thank You!",
        "Questions & Discussion"
    )
    
    # Save
    output_path = os.path.join(os.path.dirname(__file__), "..", "Division-Level-Access-Presentation.pptx")
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ python-pptx not installed. Run: pip install python-pptx")
        print("\nAlternatively, use the markdown file 'Division-Level-Access-Feature.pptx.md'")
        print("and convert it using a markdown-to-pptx tool or copy to PowerPoint manually.")
