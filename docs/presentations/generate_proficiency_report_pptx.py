#!/usr/bin/env python3
"""
Division Q2 Proficiency Level Report - PowerPoint Generator
Creates a comprehensive presentation for the automated proficiency reporting feature.

Requirements:
    pip install python-pptx

Usage:
    python generate_proficiency_report_pptx.py

Output:
    Division_Q2_Proficiency_Report_Automation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
BLUE_PRIMARY = RGBColor(30, 64, 175)      # #1E40AF - Deep Blue
BLUE_SECONDARY = RGBColor(59, 130, 246)   # #3B82F6 - Light Blue
GREEN_SUCCESS = RGBColor(34, 197, 94)     # #22C55E - Success Green
RED_ALERT = RGBColor(239, 68, 68)         # #EF4444 - Alert Red
ORANGE_WARNING = RGBColor(249, 115, 22)   # #F97316 - Warning Orange
YELLOW_HIGHLIGHT = RGBColor(234, 179, 8)  # #EAB308 - Highlight Yellow
GRAY_TEXT = RGBColor(71, 85, 105)         # #475569 - Text Gray
GRAY_LIGHT = RGBColor(241, 245, 249)      # #F1F5F9 - Light Gray
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)


def add_title_slide(prs, title, subtitle, footer_text="EduSync School Information System | December 2025"):
    """Add a title slide with gradient-like background"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE_PRIMARY
    shape.line.fill.background()
    
    # Decorative accent bar
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.2), prs.slide_width, Inches(0.1)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GREEN_SUCCESS
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(191, 219, 254)
    p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = footer_text
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(147, 197, 253)
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_section_slide(prs, title, subtitle=""):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE_SECONDARY
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(219, 234, 254)
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs, title, bullet_points, two_column=False):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE_PRIMARY
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    if two_column and len(bullet_points) > 4:
        mid = len(bullet_points) // 2
        left_points = bullet_points[:mid]
        right_points = bullet_points[mid:]
        
        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.3), Inches(5))
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, point in enumerate(left_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = GRAY_TEXT
            p.space_after = Pt(10)
        
        # Right column
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.3), Inches(5))
        tf = right_box.text_frame
        tf.word_wrap = True
        for i, point in enumerate(right_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = GRAY_TEXT
            p.space_after = Pt(10)
    else:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.3))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(18)
            p.font.color.rgb = GRAY_TEXT
            p.space_after = Pt(12)
    
    return slide


def add_comparison_slide(prs, title, before_title, before_points, after_title, after_points):
    """Add a before/after comparison slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE_PRIMARY
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # BEFORE column (red)
    before_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.3), Inches(4.5), Inches(5.2)
    )
    before_bg.fill.solid()
    before_bg.fill.fore_color.rgb = RGBColor(254, 226, 226)  # Light red
    before_bg.line.color.rgb = RED_ALERT
    
    before_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.1), Inches(0.5))
    tf = before_header.text_frame
    p = tf.paragraphs[0]
    p.text = f"❌ {before_title}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED_ALERT
    
    before_content = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.1), Inches(4.3))
    tf = before_content.text_frame
    tf.word_wrap = True
    for i, point in enumerate(before_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY_TEXT
        p.space_after = Pt(8)
    
    # AFTER column (green)
    after_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.3), Inches(4.5), Inches(5.2)
    )
    after_bg.fill.solid()
    after_bg.fill.fore_color.rgb = RGBColor(220, 252, 231)  # Light green
    after_bg.line.color.rgb = GREEN_SUCCESS
    
    after_header = slide.shapes.add_textbox(Inches(5.4), Inches(1.4), Inches(4.1), Inches(0.5))
    tf = after_header.text_frame
    p = tf.paragraphs[0]
    p.text = f"✅ {after_title}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN_SUCCESS
    
    after_content = slide.shapes.add_textbox(Inches(5.4), Inches(2), Inches(4.1), Inches(4.3))
    tf = after_content.text_frame
    tf.word_wrap = True
    for i, point in enumerate(after_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY_TEXT
        p.space_after = Pt(8)
    
    return slide


def add_stats_slide(prs, title, stats):
    """Add a slide with statistics boxes"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE_PRIMARY
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Stats boxes
    colors = [BLUE_PRIMARY, GREEN_SUCCESS, ORANGE_WARNING, BLUE_SECONDARY]
    box_width = Inches(2.2)
    box_height = Inches(1.8)
    start_x = Inches(0.4)
    start_y = Inches(1.5)
    gap = Inches(0.15)
    
    for i, (value, label) in enumerate(stats[:4]):
        x = start_x + (box_width + gap) * i
        
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, box_width, box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.fill.background()
        
        # Value
        val_box = slide.shapes.add_textbox(x, start_y + Inches(0.3), box_width, Inches(0.8))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        lbl_box = slide.shapes.add_textbox(x, start_y + Inches(1.1), box_width, Inches(0.6))
        tf = lbl_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(219, 234, 254)
        p.alignment = PP_ALIGN.CENTER
    
    # Additional stats if more than 4
    if len(stats) > 4:
        start_y2 = Inches(3.6)
        for i, (value, label) in enumerate(stats[4:8]):
            x = start_x + (box_width + gap) * i
            
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y2, box_width, box_height
            )
            box.fill.solid()
            box.fill.fore_color.rgb = colors[(i + 2) % len(colors)]
            box.line.fill.background()
            
            val_box = slide.shapes.add_textbox(x, start_y2 + Inches(0.3), box_width, Inches(0.8))
            tf = val_box.text_frame
            p = tf.paragraphs[0]
            p.text = value
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            
            lbl_box = slide.shapes.add_textbox(x, start_y2 + Inches(1.1), box_width, Inches(0.6))
            tf = lbl_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(219, 234, 254)
            p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_workflow_slide(prs, title, steps):
    """Add a workflow/process slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE_PRIMARY
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Process steps
    step_height = Inches(0.9)
    start_y = Inches(1.4)
    
    for i, (step_title, step_desc, time_estimate) in enumerate(steps):
        y = start_y + (step_height + Inches(0.1)) * i
        
        # Step number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.4), y + Inches(0.1), Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLUE_PRIMARY
        circle.line.fill.background()
        
        num_box = slide.shapes.add_textbox(Inches(0.4), y + Inches(0.15), Inches(0.5), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Step content
        content_box = slide.shapes.add_textbox(Inches(1.1), y, Inches(6.5), step_height)
        tf = content_box.text_frame
        p = tf.paragraphs[0]
        p.text = step_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLACK
        
        p2 = tf.add_paragraph()
        p2.text = step_desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = GRAY_TEXT
        
        # Time badge
        time_box = slide.shapes.add_textbox(Inches(7.8), y + Inches(0.15), Inches(1.5), Inches(0.4))
        tf = time_box.text_frame
        p = tf.paragraphs[0]
        p.text = time_estimate
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = GREEN_SUCCESS
        p.alignment = PP_ALIGN.RIGHT
    
    return slide


def add_table_slide(prs, title, headers, rows):
    """Add a slide with a simple table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE_PRIMARY
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Create table
    cols = len(headers)
    table_rows = len(rows) + 1
    
    table_width = Inches(9)
    table_height = Inches(0.4) * table_rows
    
    table = slide.shapes.add_table(
        table_rows, cols, 
        Inches(0.5), Inches(1.4),
        table_width, table_height
    ).table
    
    # Set column widths
    col_width = table_width / cols
    for i in range(cols):
        table.columns[i].width = int(col_width)
    
    # Header row
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_PRIMARY
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(12)
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.font.color.rgb = GRAY_TEXT
            para.alignment = PP_ALIGN.CENTER
            
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GRAY_LIGHT
    
    return slide


def add_quote_slide(prs, quote, author, role):
    """Add a testimonial/quote slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRAY_LIGHT
    shape.line.fill.background()
    
    # Quote mark
    quote_mark = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(1), Inches(1))
    tf = quote_mark.text_frame
    p = tf.paragraphs[0]
    p.text = '"'
    p.font.size = Pt(100)
    p.font.color.rgb = RGBColor(191, 219, 254)
    
    # Quote text
    quote_box = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(8), Inches(2.5))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'"{quote}"'
    p.font.size = Pt(24)
    p.font.italic = True
    p.font.color.rgb = GRAY_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    # Author
    author_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.8))
    tf = author_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"— {author}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = role
    p2.font.size = Pt(14)
    p2.font.color.rgb = GRAY_TEXT
    p2.alignment = PP_ALIGN.CENTER
    
    return slide


def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ====================
    # SLIDE 1: Title
    # ====================
    add_title_slide(
        prs,
        "Automated Q2 Proficiency\nLevel Reporting",
        "Transforming Division-Wide Data Collection\nFrom Weeks to Minutes",
        "EduSync School Information System | Division of Mati City"
    )
    
    # ====================
    # SLIDE 2: Agenda
    # ====================
    add_content_slide(prs, "📋 Presentation Agenda", [
        "The Current Challenge — Manual Quarterly Reporting",
        "The Hidden Costs — Time, Errors, and Stress",
        "The EduSync Solution — Automated Data Flow",
        "How It Works — System Demonstration",
        "Impact Analysis — Before vs. After",
        "Implementation Roadmap — Next Steps",
        "Q&A — Open Discussion"
    ])
    
    # ====================
    # SLIDE 3: Section - The Problem
    # ====================
    add_section_slide(prs, "📊 The Current Challenge", "Understanding the Manual Process")
    
    # ====================
    # SLIDE 4: Current Workflow
    # ====================
    add_workflow_slide(prs, "Current Q2 Proficiency Report Workflow", [
        ("Teacher Level", "Each teacher computes MPS per subject manually in Excel, then submits to School Head", "⏱️ 2-3 days"),
        ("School Level", "Principal/Registrar collects from all teachers, consolidates into one school report", "⏱️ 2-3 days"),
        ("District Level", "District Supervisor collects from 5-15 schools, merges into district file", "⏱️ 3-5 days"),
        ("Division Level", "Planning Officer collects from all districts, creates final consolidated report", "⏱️ 3-5 days"),
        ("Regional Submission", "Division submits to Regional Office before deadline", "⏱️ Total: 3-5 weeks")
    ])
    
    # ====================
    # SLIDE 5: Pain Points
    # ====================
    add_content_slide(prs, "😰 Pain Points in Current Process", [
        "Manual Data Entry — High risk of typos and formula errors",
        "Multiple Excel Versions — \"Which file is the latest?\"",
        "Delayed Submissions — Cascading delays from one late teacher",
        "No Data Validation — Impossible values slip through (120%, negative %)",
        "Copy-Paste Errors — Wrong school data in wrong row",
        "Format Inconsistency — Different templates from different schools",
        "Re-computation Requests — \"Please recalculate and resubmit\"",
        "USB/File Loss — Panic when files get corrupted",
        "No Audit Trail — \"Who submitted what and when?\""
    ], two_column=True)
    
    # ====================
    # SLIDE 6: Time Lost
    # ====================
    add_stats_slide(prs, "⏱️ Time Lost Every Quarter", [
        ("4-6 hrs", "Per Teacher\n(Extra work)"),
        ("8-12 hrs", "Per Principal\n(Consolidation)"),
        ("16-24 hrs", "Per District Supervisor"),
        ("24-40 hrs", "Division Compiler"),
        ("3-5", "Weeks Total\n(Submission delay)"),
        ("4x", "Per Year\n(Quarterly)"),
        ("High", "Stress Level\n(Deadline pressure)"),
        ("Many", "Revision Requests")
    ])
    
    # ====================
    # SLIDE 7: Section - The Solution
    # ====================
    add_section_slide(prs, "✨ The EduSync Solution", "Automated Proficiency Reporting")
    
    # ====================
    # SLIDE 8: How EduSync Works
    # ====================
    add_content_slide(prs, "🔄 How EduSync Solves This", [
        "Teachers already input grades for Report Cards (SF9, Form 138)",
        "System AUTOMATICALLY calculates MPS from existing Q2 grades",
        "School-level summary is generated in real-time",
        "District aggregation happens instantly",
        "Division report is ONE CLICK away",
        "Export matches exact DepEd Excel format",
        "No extra work for teachers — Zero duplicate entry"
    ])
    
    # ====================
    # SLIDE 9: Data Flow
    # ====================
    add_workflow_slide(prs, "EduSync Automated Data Flow", [
        ("Grades Already in System", "Teachers enter Q2 grades for individual students (for report cards)", "✅ Already done"),
        ("Auto-Compute MPS", "System calculates Mean Percentage Score per subject per grade level", "🤖 Instant"),
        ("School Summary", "Automatically aggregated for each school", "📊 Real-time"),
        ("District/Division Roll-up", "Consolidated across all schools instantly", "⚡ 1 second"),
        ("Download Excel", "Export in exact DepEd format with matching headers", "📥 1 click")
    ])
    
    # ====================
    # SLIDE 10: Report Structure
    # ====================
    add_table_slide(prs, "📋 Automated Report Structure", 
        ["Grade/Subject", "Data Captured", "Calculation"],
        [
            ("Kindergarten", "A/B/C Proficiency Levels", "% at Beginning, Developing, Consistent"),
            ("Grade 1-2 Language", "Quarterly Grades", "% with 75%+ MPS, Mean Score"),
            ("Grade 3 Mother Tongue", "Quarterly Grades", "% with 75%+ MPS, Mean Score"),
            ("Grade 1-3 Reading & Literacy", "ELLN Assessment", "% with 75%+ MPS, Mean Score"),
            ("Grade 4-6 English", "Quarterly Grades", "% with 75%+ MPS, Mean Score"),
            ("All Subjects", "Per School/District", "Grouped by District")
        ]
    )
    
    # ====================
    # SLIDE 11: Comparison
    # ====================
    add_comparison_slide(prs, "⚖️ Manual vs. Automated",
        "BEFORE (Manual)",
        [
            "3-5 weeks to compile Division report",
            "40+ hours of manual consolidation",
            "High risk of calculation errors",
            "Multiple Excel file versions",
            "Stressful deadline crunch",
            "Frequent revision requests",
            "No real-time visibility",
            "Data is 2-3 weeks old"
        ],
        "AFTER (EduSync)",
        [
            "5 minutes to generate report",
            "0 hours of manual consolidation",
            "100% accurate calculations",
            "Single source of truth",
            "Stress-free submissions",
            "Get it right the first time",
            "Real-time dashboard access",
            "Data is always current"
        ]
    )
    
    # ====================
    # SLIDE 12: Impact Stats
    # ====================
    add_stats_slide(prs, "📈 Measurable Impact", [
        ("99.9%", "Faster Report\nGeneration"),
        ("100%", "Error\nElimination"),
        ("0 hrs", "Extra Teacher\nWork"),
        ("95%", "Time Saved\n(All Levels)"),
    ])
    
    # ====================
    # SLIDE 13: Quote - Teacher
    # ====================
    add_quote_slide(prs,
        "Wala na kaming extra Excel. Yung grades na ini-input namin for report card, automatic na napupunta sa Division report.",
        "Sample Teacher Feedback",
        "Elementary School Teacher"
    )
    
    # ====================
    # SLIDE 14: Quote - Principal
    # ====================
    add_quote_slide(prs,
        "Hindi na kami nag-coconsolidate. Click lang, tapos na. Mas maaga pa kami makakapag-submit.",
        "Sample Principal Feedback",
        "Elementary School Principal"
    )
    
    # ====================
    # SLIDE 15: Quote - Division
    # ====================
    add_quote_slide(prs,
        "Instant report. No more encoding errors. No more 'Sir, mali pala yung data namin, palitan.'",
        "Sample Division Feedback",
        "Division Planning Officer"
    )
    
    # ====================
    # SLIDE 16: Section - Demo
    # ====================
    add_section_slide(prs, "🖥️ Live Demonstration", "See It In Action")
    
    # ====================
    # SLIDE 17: Demo Steps
    # ====================
    add_content_slide(prs, "Demo Walkthrough", [
        "1. Login to EduSync as Division User",
        "2. Navigate to Reports → Q2 Proficiency Level",
        "3. Select School Year: 2025-2026",
        "4. Select Quarter: Q2 (2nd Quarter)",
        "5. View real-time data from all schools",
        "6. Filter by District if needed",
        "7. Click 'Download Excel' → Done!"
    ])
    
    # ====================
    # SLIDE 18: Section - Implementation
    # ====================
    add_section_slide(prs, "🚀 Implementation Roadmap", "How Do We Get There?")
    
    # ====================
    # SLIDE 19: Prerequisites
    # ====================
    add_content_slide(prs, "✅ Prerequisites for Automation", [
        "Schools must use EduSync for Grade Entry — Teachers input Q1-Q4 grades in system",
        "Complete Student Enrollment Data — Students properly assigned to sections",
        "Learning Areas Configured — Subjects set up per grade level",
        "Division Account Setup — Division users with proper access",
        "Network/Internet Access — For syncing data to central system"
    ])
    
    # ====================
    # SLIDE 20: Rollout Plan
    # ====================
    add_table_slide(prs, "📅 Rollout Plan",
        ["Phase", "Timeline", "Schools", "Outcome"],
        [
            ("Pilot", "Jan 2026", "5 schools", "Validate report accuracy"),
            ("District 1", "Feb 2026", "All Central", "First district automated"),
            ("District 2-3", "Mar 2026", "North + South", "Expand coverage"),
            ("Full Division", "Apr 2026", "All schools", "100% automated Q3 report"),
        ]
    )
    
    # ====================
    # SLIDE 21: Support
    # ====================
    add_content_slide(prs, "🤝 Support & Training", [
        "Teacher Training — 1-hour orientation on grade entry",
        "Principal Training — Dashboard overview and report generation",
        "District Supervisor Access — Real-time monitoring dashboard",
        "Division Admin Training — Full report generation and export",
        "Ongoing Support — Help desk, documentation, video tutorials",
        "Data Migration — Assistance with historical data if needed"
    ])
    
    # ====================
    # SLIDE 22: Section - Q&A
    # ====================
    add_section_slide(prs, "❓ Questions & Discussion", "We're Here to Help")
    
    # ====================
    # SLIDE 23: Contact
    # ====================
    add_content_slide(prs, "📞 Get Started Today", [
        "Schedule a Demo — See the full system in action",
        "Pilot Program — Start with a few schools, prove the value",
        "Training Sessions — We train your teachers and staff",
        "Deployment Support — Full assistance during rollout",
        "",
        "Contact: EduSync Implementation Team",
        "Email: support@edusync.ph",
        "Website: www.edusync.ph"
    ])
    
    # ====================
    # SLIDE 24: Thank You
    # ====================
    add_title_slide(
        prs,
        "Thank You!",
        "Let's Transform Education Data Together\n\n📧 support@edusync.ph",
        "EduSync School Information System | Division of Mati City"
    )
    
    # Save presentation
    output_path = os.path.join(os.path.dirname(__file__), "Division_Q2_Proficiency_Report_Automation.pptx")
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
